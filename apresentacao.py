from pptx import Presentation
from pptx.util import Inches, Pt

# Cria uma apresentação
prs = Presentation()

# Função para adicionar slides com texto
def add_bullet_slide(title_text, bullet_points, layout_index=1):
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # Limpar o texto anterior

    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(24)  # Tamanho da fonte em pontos

    return slide

# Slide de título
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Big Data na Logística"
slide.placeholders[1].text = "O poder dos dados na eficiência dos entregadores\nSeu Nome - Data"

# Adicionar slides com tópicos
bullet_points = [
    "Grandes volumes de dados que exigem métodos específicos para armazenamento e análise",
    "As 5 V's: Volume, Velocidade, Variedade, Veracidade, Valor"
]
add_bullet_slide("O que é Big Data", bullet_points)

bullet_points = [
    "Otimizar rotas e reduzir custos",
    "Monitoramento em tempo real: controle de frotas e rastreamento",
    "Previsão de demandas: antecipação de problemas e necessidades",
    "Melhoria no atendimento: ajustes operacionais com base nos dados"
]
add_bullet_slide("Por que o Big Data é Importante na Logística?", bullet_points)

bullet_points = [
    "Empresas: Amazon, Mercado Livre, Correios",
    "Roteirização inteligente com dados (GPS, otimização de rotas)",
    "Centros de distribuição automatizados",
]
add_bullet_slide("Exemplos Reais no Mercado", bullet_points)

bullet_points = [
    "Demonstração prática de Big Data aplicada à logística",
    "Geração de dados de entregas e feedbacks de clientes",
    "Dashboard interativo com visualizações",
]
add_bullet_slide("Nosso sistema: LogiData", bullet_points)

bullet_points = [
    "Geração de dados: entregas (cidades, tempo, custo, status) e feedbacks",
    "Análise quantitativa: gráficos de entregas, custos e atrasos",
    "Análise qualitativa: processamento de feedbacks e análise de sentimentos",
    "Visualizações interativas: dashboards e nuvem de palavras"
]
add_bullet_slide("Como Funciona o Sistema", bullet_points)

bullet_points = [
    "Gráfico de entregas por cidade",
    "Custo médio por transportadora",
    "Indicadores de desempenho: número de entregas, atrasos, custos",
    "Nuvem de palavras e análise de sentimentos dos feedbacks"
]
add_bullet_slide("Resultados e Visualizações", bullet_points)

bullet_points = [
    "Big Data proporciona vantagem competitiva na logística",
    "Melhora na tomada de decisão e otimização de operações",
    "Integração de análises quantitativas e qualitativas gera insights valiosos",
    "Futuro: avanço com machine learning e análises preditivas"
]
add_bullet_slide("Conclusão", bullet_points)

bullet_points = [
    "Obrigado!",
    "Perguntas",
    "Seu Nome - contato@seuemail.com",
]
add_bullet_slide("Perguntas & Contato", bullet_points)

# Salvar a apresentação
arquivo = "apresentacao_big_data.pptx"
prs.save(arquivo)
print(f"Apresentação salva como '{arquivo}'")